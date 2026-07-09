from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Palette (matches my earlier decks) ────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DEEP_RED   = RGBColor(0x8B, 0x00, 0x00)
BLACK      = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_BODY  = RGBColor(0x33, 0x33, 0x33)
TEXT_LIGHT = RGBColor(0x66, 0x66, 0x66)
RULE_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)

FONT_TITLE = "Arial"
FONT_BODY  = "Arial"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height
TOTAL_SLIDES = 8

# ── Helpers ────────────────────────────────────────────────────────────────
def white_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

def rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def line(slide, left, top, width, color=DEEP_RED, thickness=Pt(2)):
    return rect(slide, left, top, width, thickness, color)

def txt(slide, left, top, width, height, text, size=18, color=TEXT_BODY,
        bold=False, align=PP_ALIGN.LEFT, font=FONT_BODY, spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    p.line_spacing = Pt(size * spacing)
    return box

def bullets(slide, left, top, width, height, items, size=15, color=TEXT_BODY,
            spacing=1.6, bold_prefix=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = Pt(size * spacing)
        p.space_after = Pt(6)
        if bold_prefix and ": " in item:
            prefix, rest = item.split(": ", 1)
            r1 = p.add_run(); r1.text = "•   " + prefix + ": "
            r1.font.size = Pt(size); r1.font.color.rgb = BLACK; r1.font.bold = True; r1.font.name = FONT_BODY
            r2 = p.add_run(); r2.text = rest
            r2.font.size = Pt(size); r2.font.color.rgb = color; r2.font.bold = False; r2.font.name = FONT_BODY
        else:
            r = p.add_run(); r.text = "•   " + item
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = False; r.font.name = FONT_BODY
    return box

def slide_number(slide, num):
    txt(slide, W - Inches(1.3), H - Inches(0.5), Inches(1), Inches(0.3),
        f"{num} / {TOTAL_SLIDES}", size=10, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)

def section_title(slide, title, subtitle=None):
    line(slide, Inches(0.9), Inches(0.85), Inches(0.9), color=DEEP_RED, thickness=Pt(3))
    txt(slide, Inches(0.9), Inches(1.05), Inches(11), Inches(0.7),
        title, size=36, color=DEEP_RED, bold=True, font=FONT_TITLE)
    if subtitle:
        txt(slide, Inches(0.9), Inches(1.7), Inches(11), Inches(0.45),
            subtitle, size=16, color=TEXT_LIGHT)

def base():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    white_bg(s)
    return s

def col_header(slide, x, y, w, text):
    txt(slide, x, y, w, Inches(0.4), text, size=18, color=DEEP_RED, bold=True)
    line(slide, x, y + Inches(0.45), Inches(2.5), color=DEEP_RED, thickness=Pt(1.5))

# ═══ SLIDE 1 — TITLE ═════════════════════════════════════════════════════
s = base()
txt(s, Inches(0), Inches(2.0), W, Inches(1.0),
    "INSIDER LEDGER", size=54, color=DEEP_RED, bold=True, align=PP_ALIGN.CENTER, font=FONT_TITLE)
line(s, Inches(5.4), Inches(3.15), Inches(2.5), color=DEEP_RED, thickness=Pt(2.5))
txt(s, Inches(0), Inches(3.5), W, Inches(0.6),
    "Who's Buying and Selling Their Own Stock?", size=22, color=TEXT_BODY, align=PP_ALIGN.CENTER)
txt(s, Inches(0), Inches(4.2), W, Inches(0.6),
    "A searchable record of corporate insider trades, from SEC filings and daily stock prices.",
    size=16, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
txt(s, Inches(0), H - Inches(0.75), W, Inches(0.4),
    "Dimuthu Attanayake  |  Data and Databases  |  July 9", size=12, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
slide_number(s, 1)

# ═══ SLIDE 2 — THE TOPIC ═════════════════════════════════════════════════
s = base()
section_title(s, "The Topic", "What the data is, and why it matters")
bullets(s, Inches(1.1), Inches(2.5), Inches(11.2), Inches(3),
    [
        "Directors, officers, and big shareholders must file a Form 4 with the SEC within two business days of buying or selling their own company's stock.",
        "It is one of the few public windows into what the people who know a company best are doing with their money.",
        "The filings are public, but scattered and full of jargon, so a normal reader can't really use them.",
    ], size=16, spacing=1.7)
rect(s, Inches(1.0), Inches(5.35), Inches(0.05), Inches(1.05), DEEP_RED)
txt(s, Inches(1.35), Inches(5.35), Inches(10.8), Inches(0.4),
    "My question", size=15, color=DEEP_RED, bold=True)
txt(s, Inches(1.35), Inches(5.75), Inches(10.8), Inches(0.7),
    "Who is trading their own stock right now, are they buying or selling, and are they trading above or below the market price?",
    size=17, color=BLACK, spacing=1.3)
slide_number(s, 2)

# ═══ SLIDE 3 — THE DATA (two columns) ════════════════════════════════════
s = base()
section_title(s, "The Data: Two Sources", "One scraped, one from an API, merged on ticker and date")
col_header(s, Inches(1.1), Inches(2.5), Inches(5), "Source 1: SEC EDGAR (scraped)")
bullets(s, Inches(1.1), Inches(3.15), Inches(5.1), Inches(3),
    [
        "The latest Form 4 filings feed.",
        "Gives the insider, the company, the role, and the actual trades: shares, price, buy or sell.",
        "The primary record, collected daily.",
    ], size=14, spacing=1.7)
rect(s, Inches(6.55), Inches(2.5), Inches(0.02), Inches(4.0), RULE_GRAY)
col_header(s, Inches(7.0), Inches(2.5), Inches(5), "Source 2: Yahoo Finance (API)")
bullets(s, Inches(7.0), Inches(3.15), Inches(5.2), Inches(3),
    [
        "The daily closing price for each company.",
        "Answers the part Form 4 can't: did the insider trade above or below what the public paid that day?",
        "The context that makes one trade meaningful.",
    ], size=14, spacing=1.7)
slide_number(s, 3)

# ═══ SLIDE 4 — THE PROCESS ═══════════════════════════════════════════════
s = base()
section_title(s, "The Process", "A pipeline that updates itself")
bullets(s, Inches(1.1), Inches(2.5), Inches(11.2), Inches(3.5),
    [
        "Scrape: pull the latest 100 Form 4 filings, then follow each to its XML for the full trade details and footnote text.",
        "Detect and log: each filing gets a hash, so I can tell new filings from old ones and skip re-scraping. Every run writes a change log and an error log.",
        "Automate: a GitHub Action runs the scrape every day and saves new filings, so the data grows on its own. It went from 100 to about 300 rows in three days.",
        "Merge: I then pull stock prices for each ticker and join them to the trades.",
    ], size=16, spacing=1.8, bold_prefix=True)
slide_number(s, 4)

# ═══ SLIDE 5 — CHALLENGES (cards) ════════════════════════════════════════
s = base()
section_title(s, "Challenges So Far", "The main things that fought back")
cards = [
    ("The same filing appears twice",
     "Each Form 4 is listed once under the insider and once under the company, so 100 rows was really 42 filings.",
     "Fix: dedupe on the accession number."),
    ("The data is deeply nested",
     "The real detail lives in XML, with a different shape for different trade types.",
     "Fix: flatten it into clean rows, keeping only the fields that exist."),
    ("Prices don't always match",
     "Weekend trades have no closing price, and small or foreign tickers aren't on Yahoo.",
     "Result: about 10% of trades get no market price, left blank."),
]
cy, ch, gap = Inches(2.5), Inches(1.35), Inches(0.2)
for i, (title, body, fix) in enumerate(cards):
    y = cy + i * (ch + gap)
    rect(s, Inches(1.0), y, Inches(0.05), ch, DEEP_RED)
    txt(s, Inches(1.35), y + Inches(0.03), Inches(11), Inches(0.35), title, size=18, color=BLACK, bold=True)
    txt(s, Inches(1.35), y + Inches(0.43), Inches(11), Inches(0.5), body, size=14, color=TEXT_BODY)
    txt(s, Inches(1.35), y + Inches(0.95), Inches(11), Inches(0.35), fix, size=13, color=TEXT_LIGHT)
slide_number(s, 5)

# ═══ SLIDE 6 — THE SCHEMA (two columns) ══════════════════════════════════
s = base()
section_title(s, "The Schema", "Two tables, joined. One filing has many trades.")
col_header(s, Inches(1.1), Inches(2.5), Inches(5), "filings  (one row per filing)")
bullets(s, Inches(1.1), Inches(3.15), Inches(5.1), Inches(3),
    ["accession (id)", "insider name", "company, ticker", "relationship", "filing date", "footnotes, SEC link"],
    size=14, spacing=1.6)
rect(s, Inches(6.55), Inches(2.5), Inches(0.02), Inches(3.6), RULE_GRAY)
col_header(s, Inches(7.0), Inches(2.5), Inches(5.2), "trades  (one row per transaction)")
bullets(s, Inches(7.0), Inches(3.15), Inches(5.3), Inches(3),
    ["accession (link to filing)", "transaction date", "code (buy or sell)", "shares, price, value",
     "market close  (from source 2)", "percent vs market  (from source 2)"],
    size=14, spacing=1.6)
txt(s, Inches(1.1), Inches(6.15), Inches(11.2), Inches(0.8),
    "One row per transaction is the searchable unit, because that is how a reader thinks: show me the buys. "
    "The two price fields join in on ticker and transaction date.",
    size=13, color=TEXT_LIGHT, spacing=1.4)
slide_number(s, 6)

# ═══ SLIDE 7 — THE PUBLIC APP (two columns) ══════════════════════════════
s = base()
section_title(s, "The Public App", "Made searchable for someone who has never read a filing")
col_header(s, Inches(1.1), Inches(2.5), Inches(5), "What users can do")
bullets(s, Inches(1.1), Inches(3.15), Inches(5.1), Inches(3),
    ["Search by insider, company, ticker, or role.", "Filter by buy vs sell.",
     "Page through results.", "Export the results to CSV.",
     "Built with Flask and SQLite, hosted on Render."],
    size=14, spacing=1.6)
rect(s, Inches(6.55), Inches(2.5), Inches(0.02), Inches(3.4), RULE_GRAY)
col_header(s, Inches(7.0), Inches(2.5), Inches(5.2), "Example searches")
bullets(s, Inches(7.0), Inches(3.15), Inches(5.3), Inches(3),
    ["Example 1: every open-market purchase by a director. Selling is routine, but an insider buying with their own cash is the rare, interesting event.",
     "Example 2: all trades for one ticker, showing whether each was above or below market."],
    size=14, spacing=1.7, bold_prefix=True)
txt(s, Inches(1.1), Inches(6.4), Inches(11), Inches(0.4),
    "(Wireframe sketch to be added here.)", size=12, color=TEXT_LIGHT)
slide_number(s, 7)

# ═══ SLIDE 8 — METHODOLOGY, HORIZON & QUESTIONS (two columns) ════════════
s = base()
section_title(s, "Methodology, Horizon & Questions")
col_header(s, Inches(1.1), Inches(2.3), Inches(5), "The data, honestly")
bullets(s, Inches(1.1), Inches(2.95), Inches(5.1), Inches(3.5),
    ["Reveals: a near real-time, verifiable record. Every row links back to the original SEC filing.",
     "Missing: recent filings only, no long history. Mostly selling and grants; buys are rare. It shows what, not why.",
     "Ethics: public data, but it aggregates named people's finances. A trade is not proof of intent."],
    size=13.5, spacing=1.6, bold_prefix=True)
rect(s, Inches(6.55), Inches(2.3), Inches(0.02), Inches(4.1), RULE_GRAY)
col_header(s, Inches(7.0), Inches(2.3), Inches(5.2), "Where next, and asks")
bullets(s, Inches(7.0), Inches(2.95), Inches(5.3), Inches(1.8),
    ["Ship the app on Render, finalize the tables, add filters, write the full methodology."],
    size=13.5, spacing=1.6)
txt(s, Inches(7.0), Inches(4.05), Inches(5.2), Inches(0.4), "Questions for you", size=14, color=DEEP_RED, bold=True)
bullets(s, Inches(7.0), Inches(4.5), Inches(5.3), Inches(2),
    ["Is one row per transaction the clearest unit, or should I group by filing or insider?",
     "How do I show 'traded 2% below market' as useful context without implying a signal that isn't there?"],
    size=13.5, spacing=1.6)
slide_number(s, 8)

out = "/Users/dimuthuattanayake/Documents/Data and Databases/Bill Scraper/Insider_Ledger_Slides.pptx"
prs.save(out)
print("saved:", out, "with", len(prs.slides._sldIdLst), "slides")
